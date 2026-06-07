"""
生成 5 頁 PowerPoint 簡報
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 色票 ──────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)   # 深藍（標題背景）
TEAL   = RGBColor(0x00, 0x87, 0x87)   # 青（強調色）
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF6, 0xF8)   # 淺灰（表格底色）
DGRAY  = RGBColor(0x44, 0x44, 0x44)   # 深灰（內文）
ORANGE = RGBColor(0xE8, 0x6A, 0x1A)   # 橘（地雷標示）
GREEN  = RGBColor(0x27, 0xAE, 0x60)   # 綠（成功標示）

W = Inches(13.33)   # 寬螢幕 16:9
H = Inches(7.5)

# ── helpers ───────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank)

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=DGRAY,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=None):
    """Deep-navy top bar with title."""
    rect(slide, 0, 0, W, Inches(1.55), fill=NAVY)
    txbox(slide, title,
          Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
          size=32, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.5),
              size=16, bold=False, color=RGBColor(0xAA,0xCC,0xDD))

def accent_bar(slide, y=Inches(1.55)):
    """Thin teal rule under header."""
    r = rect(slide, 0, y, W, Pt(4), fill=TEAL)

def add_table(slide, headers, rows, l, t, w, h,
              hdr_fill=NAVY, hdr_color=WHITE,
              row_fills=(LGRAY, WHITE)):
    from pptx.util import Inches, Pt
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols, l, t, w, h).table

    # Column widths – equal split
    col_w = w // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    def set_cell(cell, text, fill=None, fcolor=DGRAY,
                 bold=False, sz=14, align=PP_ALIGN.LEFT):
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(sz)
        run.font.bold  = bold
        run.font.color.rgb = fcolor

    for ci, h_text in enumerate(headers):
        set_cell(tbl.cell(0, ci), h_text,
                 fill=hdr_fill, fcolor=hdr_color,
                 bold=True, sz=14, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        fill = row_fills[ri % 2]
        for ci, val in enumerate(row):
            set_cell(tbl.cell(ri+1, ci), val, fill=fill, sz=13)

    return tbl

def bullet_box(slide, items, l, t, w, h,
               size=16, color=DGRAY, bullet="▸ "):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color

# ══════════════════════════════════════════════════════════
# SLIDE 1 — 封面 / 任務理解
# ══════════════════════════════════════════════════════════
def slide1(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, W, H, fill=NAVY)
    # teal accent strip
    rect(s, 0, Inches(4.6), W, Pt(5), fill=TEAL)

    txbox(s, "技能孵化製作報告",
          Inches(1), Inches(1.2), Inches(11), Inches(1.2),
          size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, "卡米爾 AI 專案經理面試實作",
          Inches(1), Inches(2.5), Inches(11), Inches(0.8),
          size=22, color=RGBColor(0xAA,0xCC,0xDD), align=PP_ALIGN.CENTER)
    txbox(s, "柏融 Ho Po-Jung　｜　2026.06.08",
          Inches(1), Inches(3.2), Inches(11), Inches(0.6),
          size=18, color=RGBColor(0x88,0xAA,0xCC), align=PP_ALIGN.CENTER)

    # 四個資料 badge
    badges = [
        ("📄 PDF", "偵查卷宗 9頁"),
        ("🔊 WAV", "廣播60分鐘"),
        ("📊 XLS", "人口統計25縣市"),
        ("📦 ZIP", "100份公文"),
    ]
    bw = Inches(2.8)
    bh = Inches(1.1)
    gap = Inches(0.35)
    total = 4 * bw + 3 * gap
    sx = (W - total) / 2
    for i, (icon_label, desc) in enumerate(badges):
        x = sx + i * (bw + gap)
        y = Inches(5.0)
        rect(s, x, y, bw, bh,
             fill=RGBColor(0x22, 0x44, 0x66), line=TEAL)
        txbox(s, icon_label,
              x + Inches(0.1), y + Inches(0.05), bw - Inches(0.2), Inches(0.5),
              size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        txbox(s, desc,
              x + Inches(0.1), y + Inches(0.55), bw - Inches(0.2), Inches(0.45),
              size=13, color=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)

    txbox(s, "Harness：Hermes Agent  ×  Gemma 4 31B（NCHC API）",
          Inches(1), Inches(6.6), Inches(11), Inches(0.5),
          size=14, color=TEAL, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — 技能包架構
# ══════════════════════════════════════════════════════════
def slide2(prs):
    s = blank_slide(prs)
    header_bar(s, "技能包架構", "什麼是「技能包」？如何讓 AI 讀懂原始資料？")
    accent_bar(s)

    # 定義說明
    rect(s, Inches(0.4), Inches(1.75), Inches(12.5), Inches(0.75),
         fill=RGBColor(0xE8,0xF4,0xF8))
    txbox(s, "技能包 = 把原始資料轉成 AI 看得懂的格式，讓 Gemma 在問答時能查找、引用、精確回答。",
          Inches(0.55), Inches(1.8), Inches(12.2), Inches(0.65),
          size=16, color=NAVY)

    # 流程圖
    flow_items = [
        ("原始資料\nPDF/WAV/XLS/ZIP", NAVY),
        ("前處理腳本\nOCR / Whisper\n/ xlrd / zipfile", TEAL),
        ("references/\nMarkdown 純文字", RGBColor(0x2E,0x7D,0x32)),
        ("NCHC Gemma\n4 31B API", RGBColor(0x6A,0x1A,0x8A)),
        ("✓ 正確答案", RGBColor(0xE8,0x6A,0x1A)),
    ]
    fw = Inches(2.2)
    fh = Inches(1.0)
    fgap = Inches(0.25)
    fy = Inches(2.7)
    total_fw = len(flow_items)*fw + (len(flow_items)-1)*fgap
    fx_start = (W - total_fw) / 2
    for i, (label, clr) in enumerate(flow_items):
        x = fx_start + i*(fw+fgap)
        rect(s, x, fy, fw, fh, fill=clr)
        txbox(s, label, x+Inches(0.05), fy+Inches(0.05),
              fw-Inches(0.1), fh-Inches(0.1),
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(flow_items)-1:
            txbox(s, "→", x+fw+Inches(0.02), fy+Inches(0.3),
                  fgap, fh-Inches(0.6),
                  size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 目錄樹
    tree = (
        "skills/cameo-interview/\n"
        "  ├── SKILL.md              ← 技能說明書（規格＋呼叫方式）\n"
        "  ├── references/\n"
        "  │   ├── pdf/              ← OCR 後卷宗全文\n"
        "  │   ├── wav/              ← Whisper 逐字稿（時間戳記）\n"
        "  │   ├── xls/              ← 縣市人口結構化表格\n"
        "  │   └── zip/              ← 100 份公文全文＋索引\n"
        "  └── qa/\n"
        "      └── all_qa.json       ← 20 題自驗 Q&A"
    )
    rect(s, Inches(0.4), Inches(3.9), Inches(12.5), Inches(2.9),
         fill=RGBColor(0x1A,0x1A,0x2E))
    txbox(s, tree, Inches(0.55), Inches(4.0), Inches(12.2), Inches(2.7),
          size=13, color=RGBColor(0x88,0xFF,0x88))

# ══════════════════════════════════════════════════════════
# SLIDE 3 — 前處理方法
# ══════════════════════════════════════════════════════════
def slide3(prs):
    s = blank_slide(prs)
    header_bar(s, "前處理：4 種資料，4 種方法", "每種格式都有獨特挑戰")
    accent_bar(s)

    cards = [
        ("📄 PDF｜偵查卷宗",
         NAVY,
         ["挑戰：每頁都是掃描圖片，直接讀文字得到 0 字",
          "方法：PyMuPDF 渲染 300 DPI → Tesseract OCR（繁中＋英）",
          "結果：9 頁全辨識，含案號、被告、犯罪事實、筆錄"],
         "🔑 key: 圖片 PDF 一定要 OCR"),
        ("🔊 WAV｜廣播音頻",
         TEAL,
         ["挑戰：IMA ADPCM 格式、台語混雜國語、60 分鐘",
          "方法：soundfile 讀取 → Whisper medium 轉錄（本地 GPU）",
          "結果：1,551 段逐字稿，含分鐘:秒時間戳記"],
         "🔑 key: 雲端封鎖 HuggingFace，改本地跑"),
        ("📊 XLS｜人口統計",
         RGBColor(0x2E,0x7D,0x32),
         ["挑戰：12 工作表、合併儲存格、pandas 讀出全是 Unnamed",
          "方法：改用 xlrd 逐格讀取，手動重建縣市×性別×年齡結構",
          "結果：25 縣市，完整年齡組別，加入 65 歲以上老年估算"],
         "🔑 key: 合併儲存格要繞過 pandas"),
        ("📦 ZIP｜公文集",
         RGBColor(0x6A,0x1A,0x8A),
         ["挑戰：100 份 PDF 公文，檔名中文編碼問題",
          "方法：zipfile 解壓 → PyMuPDF 萃取 → 逐份 Markdown",
          "結果：201 個輸出檔（含索引），全文可搜尋"],
         "🔑 key: 編碼問題不影響文件內容"),
    ]

    cw = Inches(5.9)
    ch = Inches(2.4)
    positions = [
        (Inches(0.3), Inches(1.7)),
        (Inches(7.0), Inches(1.7)),
        (Inches(0.3), Inches(4.3)),
        (Inches(7.0), Inches(4.3)),
    ]
    for (cx, cy), (title, color, bullets, key) in zip(positions, cards):
        rect(s, cx, cy, cw, ch, fill=LGRAY)
        rect(s, cx, cy, cw, Inches(0.42), fill=color)
        txbox(s, title, cx+Inches(0.1), cy+Inches(0.04),
              cw-Inches(0.2), Inches(0.36),
              size=15, bold=True, color=WHITE)
        for bi, b in enumerate(bullets):
            txbox(s, "• " + b,
                  cx+Inches(0.12), cy+Inches(0.46)+bi*Inches(0.48),
                  cw-Inches(0.24), Inches(0.45),
                  size=12, color=DGRAY)
        rect(s, cx, cy+Inches(1.9), cw, Inches(0.45),
             fill=RGBColor(0xFF,0xF3,0xE0))
        txbox(s, key, cx+Inches(0.1), cy+Inches(1.93),
              cw-Inches(0.2), Inches(0.38),
              size=12, bold=True, color=ORANGE)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — 自驗 20 題：成功、地雷與部署測試
# ══════════════════════════════════════════════════════════
def slide4(prs):
    s = blank_slide(prs)
    header_bar(s, "自驗 20 題：成功、地雷與部署測試", "前處理 + Railway 部署 + 向量搜尋調優")
    accent_bar(s)

    # 成功範例表（左上）
    txbox(s, "✅ 成功範例",
          Inches(0.4), Inches(1.75), Inches(5.8), Inches(0.35),
          size=13, bold=True, color=GREEN)
    add_table(s,
        ["來源", "問題", "答案節錄"],
        [
            ["PDF", "被告帳戶如何遭利用？", "ATM 分4筆提領12萬"],
            ["XLS", "新北市總人口？",        "4,046,564 人"],
            ["ZIP", "公文001主旨？",          "核定補助案計 33 件"],
            ["WAV", "余勝福參選哪裡？",       "大寮-林門市議員"],
        ],
        l=Inches(0.4), t=Inches(2.12),
        w=Inches(5.8), h=Inches(1.9),
        hdr_fill=GREEN,
    )

    # 前處理地雷（右上）
    txbox(s, "⚠️ 前處理地雷",
          Inches(6.6), Inches(1.75), Inches(6.3), Inches(0.35),
          size=13, bold=True, color=ORANGE)
    add_table(s,
        ["地雷", "原因", "解法"],
        [
            ["PDF 無文字",   "掃描圖片格式",           "Tesseract OCR"],
            ["XLS 欄名亂碼", "合併儲存格→Unnamed",     "xlrd 逐格讀取"],
            ["WAV 解碼失敗", "IMA ADPCM 格式",          "換 soundfile"],
            ["模型下載被擋", "雲端封鎖 HuggingFace",    "本地執行 Whisper"],
        ],
        l=Inches(6.6), t=Inches(2.12),
        w=Inches(6.3), h=Inches(1.9),
        hdr_fill=ORANGE,
    )

    # 部署測試地雷（下方全寬）
    txbox(s, "🚀 部署測試（Railway）發現的問題",
          Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.35),
          size=13, bold=True, color=RGBColor(0x6A,0x1A,0x8A))
    add_table(s,
        ["問題", "根本原因", "解法"],
        [
            ["Build 失敗：無法偵測語言",   "根目錄無 requirements.txt",    "根目錄加 requirements.txt"],
            ["Build 失敗：無 start cmd",   "main.py 不在根目錄",           "根目錄加 main.py 引入 app"],
            ["向量搜尋未啟用（fulltext）", "chromadb 未列入 requirements", "補加 chromadb==1.5.9"],
            ["ZIP 找不到特定公文編號",      "TOP_K=5 + 100份文件相似度接近", "ZIP TOP_K: 5 → 15"],
            ["429 Too Many Requests",      "連續請求觸發 NCHC 流量限制",   "Retry + backoff（10/20/40s）"],
            ["中文答題準確率偏低",          "英文 embedding 模型",          "換多語言 MiniLM-L12-v2"],
        ],
        l=Inches(0.4), t=Inches(4.52),
        w=Inches(12.5), h=Inches(2.65),
        hdr_fill=RGBColor(0x6A,0x1A,0x8A),
    )

# ══════════════════════════════════════════════════════════
# SLIDE 5 — 時間軸與可延伸方向
# ══════════════════════════════════════════════════════════
def slide5(prs):
    s = blank_slide(prs)
    header_bar(s, "時間軸與可延伸方向", "總製作時間約 6 小時（含部署測試）")
    accent_bar(s)

    # 時間軸（左半）
    txbox(s, "⏱ 製作時間軸",
          Inches(0.4), Inches(1.75), Inches(5.5), Inches(0.4),
          size=15, bold=True, color=NAVY)

    timeline = [
        ("資料確認＋環境建置",     "~1 hr",   TEAL),
        ("PDF OCR 前處理",         "~10 min", RGBColor(0x2E,0x7D,0x32)),
        ("XLS 結構解析（含debug）", "~1 hr",   RGBColor(0x6A,0x1A,0x8A)),
        ("ZIP 批量萃取",            "~5 min",  RGBColor(0x2E,0x7D,0x32)),
        ("WAV 轉錄（Whisper）",     "~30 min", TEAL),
        ("20 題 QA 撰寫＋自驗",     "~1 hr",   NAVY),
        ("向量搜尋實作（ChromaDB）", "~1 hr",   TEAL),
        ("Railway 部署＋除錯",      "~1 hr",   ORANGE),
    ]
    ty = Inches(2.2)
    for i, (label, time, color) in enumerate(timeline):
        rect(s, Inches(0.4), ty + i*Inches(0.68),
             Inches(0.22), Inches(0.45), fill=color)
        txbox(s, label,
              Inches(0.75), ty + i*Inches(0.68) + Inches(0.04),
              Inches(3.6), Inches(0.45), size=14, color=DGRAY)
        txbox(s, time,
              Inches(4.4), ty + i*Inches(0.68) + Inches(0.04),
              Inches(1.4), Inches(0.45), size=14, bold=True, color=color,
              align=PP_ALIGN.RIGHT)

    rect(s, Inches(0.4), ty + len(timeline)*Inches(0.68),
         Inches(5.5), Pt(2), fill=NAVY)
    txbox(s, "總計：約 6 小時",
          Inches(0.4), ty + len(timeline)*Inches(0.68) + Inches(0.1),
          Inches(5.5), Inches(0.5),
          size=18, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    # 可延伸方向（右半）
    txbox(s, "🚀 可延伸方向",
          Inches(7.0), Inches(1.75), Inches(5.9), Inches(0.4),
          size=15, bold=True, color=NAVY)

    extensions = [
        ("✓ RAG 向量搜尋（已完成）",
         "ChromaDB，603 chunks，TOP_K per source，只注入相關段落",
         GREEN, RGBColor(0xE8,0xF8,0xED)),
        ("✓ Railway API 部署（已完成）",
         "/ask + /health + /sources，向量搜尋啟用，fulltext fallback",
         GREEN, RGBColor(0xE8,0xF8,0xED)),
        ("多語言 Embedding（進行中）",
         "從英文 all-MiniLM 換成 multilingual-MiniLM-L12-v2，提升中文準確率",
         ORANGE, RGBColor(0xFF,0xF3,0xE0)),
        ("Function Calling / Tool Use",
         "給 Gemma 工具（search_pdf/wav/xls/zip），自行決定呼叫哪個來源",
         TEAL, LGRAY),
        ("多文件交叉查詢",
         "同時引用 PDF 卷宗＋XLS 人口資料，回答跨來源的複合問題",
         TEAL, LGRAY),
    ]
    for i, (title, desc, title_color, bg_color) in enumerate(extensions):
        ey = Inches(2.2) + i * Inches(0.95)
        rect(s, Inches(7.0), ey, Inches(5.9), Inches(0.85),
             fill=bg_color, line=title_color)
        txbox(s, title,
              Inches(7.15), ey + Inches(0.04),
              Inches(5.6), Inches(0.36),
              size=14, bold=True, color=title_color)
        txbox(s, desc,
              Inches(7.15), ey + Inches(0.40),
              Inches(5.6), Inches(0.40),
              size=12, color=DGRAY)

    # 底部
    rect(s, 0, Inches(7.0), W, Inches(0.5), fill=NAVY)
    txbox(s, "現場歡迎提問任何技術細節　｜　repo: hopojung/Cameo　branch: claude/cameo-interview-prep-0cOgW",
          Inches(0.3), Inches(7.02), Inches(12.7), Inches(0.4),
          size=11, color=RGBColor(0xAA,0xCC,0xDD), align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    out = "slides/報告_柏融.pptx"
    prs.save(out)
    print(f"✓ 儲存至 {out}")
